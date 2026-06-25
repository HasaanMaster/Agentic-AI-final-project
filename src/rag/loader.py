"""Reads the files in knowledge_base/ into text, remembering each piece's source."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pypdf import PdfReader


@dataclass
class Document:
    """A piece of source material and where it came from."""
    text: str
    source: str
    location: str


def _shape_text(shape) -> str:
    parts: list[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs)
            if line.strip():
                parts.append(line)

    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                parts.append(row_text)

    if shape.shape_type == 6:  # a group of shapes - look inside it
        for sub_shape in shape.shapes:
            sub = _shape_text(sub_shape)
            if sub:
                parts.append(sub)

    return "\n".join(parts)


def _load_pptx(path: Path) -> list[Document]:
    docs: list[Document] = []
    presentation = Presentation(str(path))

    for slide_number, slide in enumerate(presentation.slides, start=1):
        pieces = [_shape_text(shape) for shape in slide.shapes]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                pieces.append(f"(notes) {notes.strip()}")
        text = "\n".join(p for p in pieces if p.strip())
        if text.strip():
            docs.append(Document(text=text, source=path.name, location=f"slide {slide_number}"))

    return docs


def _load_pdf(path: Path) -> list[Document]:
    docs: list[Document] = []
    reader = PdfReader(str(path))
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            docs.append(Document(text=text, source=path.name, location=f"page {page_number}"))
    return docs


def _load_txt(path: Path) -> list[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    if not text.strip():
        return []
    return [Document(text=text, source=path.name, location="full text")]


_LOADERS = {".pptx": _load_pptx, ".pdf": _load_pdf, ".txt": _load_txt}


def load_documents(knowledge_base_dir: Path) -> list[Document]:
    """Read every supported file in the knowledge base into a flat list of Documents."""
    knowledge_base_dir = Path(knowledge_base_dir)
    documents: list[Document] = []
    for path in sorted(knowledge_base_dir.iterdir()):
        loader = _LOADERS.get(path.suffix.lower())
        if loader is not None:
            documents.extend(loader(path))
    return documents
